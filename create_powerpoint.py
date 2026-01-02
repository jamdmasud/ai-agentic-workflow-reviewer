#!/usr/bin/env python3
"""
PowerPoint Generator for Agentic Workflow Reviewer Presentation
Creates a .pptx file from the structured presentation data.

Requirements:
    pip install python-pptx

Usage:
    python create_powerpoint.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json

def create_presentation():
    """Create the PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    
    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    BLUE = RGBColor(102, 126, 234)  # #667eea
    PURPLE = RGBColor(118, 75, 162)  # #764ba2
    GOLD = RGBColor(255, 215, 0)    # #FFD700
    GREEN = RGBColor(144, 238, 144)  # #90EE90
    WHITE = RGBColor(255, 255, 255)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (approximated with solid color)
    background = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Agentic Workflow Reviewer"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.color.rgb = WHITE
    title_para.font.bold = True
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Multi-Agent Workflow Analysis & Optimization"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = WHITE
    
    # Built with Kiro
    kiro_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(0.8))
    kiro_frame = kiro_box.text_frame
    kiro_frame.text = "Built with Kiro.dev"
    kiro_para = kiro_frame.paragraphs[0]
    kiro_para.alignment = PP_ALIGN.CENTER
    kiro_para.font.size = Pt(20)
    kiro_para.font.color.rgb = GOLD
    
    # Date
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "Kiro Hackathon • January 2026"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = WHITE
    
    # Slide 2: Problem Statement & Vision
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = BLUE
    bg2.line.fill.background()
    
    # Title
    title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title2_frame = title2_box.text_frame
    title2_frame.text = "Problem Statement & Vision"
    title2_para = title2_frame.paragraphs[0]
    title2_para.font.size = Pt(36)
    title2_para.font.color.rgb = GOLD
    title2_para.font.bold = True
    
    # Challenge section
    challenge_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(3))
    challenge_frame = challenge_box.text_frame
    challenge_frame.text = "The Challenge"
    challenge_para = challenge_frame.paragraphs[0]
    challenge_para.font.size = Pt(24)
    challenge_para.font.color.rgb = GREEN
    challenge_para.font.bold = True
    
    challenges = [
        "DevOps teams struggle with complex workflow optimization",
        "Manual review processes are time-consuming and error-prone",
        "Different optimization goals require different expertise",
        "No unified system for comprehensive workflow analysis"
    ]
    
    for challenge in challenges:
        p = challenge_frame.add_paragraph()
        p.text = f"• {challenge}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
    
    # Solution section
    solution_box = slide2.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6), Inches(3))
    solution_frame = solution_box.text_frame
    solution_frame.text = "Why Agents + Kiro?"
    solution_para = solution_frame.paragraphs[0]
    solution_para.font.size = Pt(24)
    solution_para.font.color.rgb = GREEN
    solution_para.font.bold = True
    
    solutions = [
        "Decomposition: Complex analysis → specialized agents",
        "Expertise: Each agent focuses on domain knowledge",
        "Scalability: Parallel processing + intelligent caching",
        "Human-in-the-loop: Interactive goal changes"
    ]
    
    for solution in solutions:
        p = solution_frame.add_paragraph()
        p.text = f"• {solution}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
    
    # Slide 3: Architecture
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = BLUE
    bg3.line.fill.background()
    
    # Title
    title3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title3_frame = title3_box.text_frame
    title3_frame.text = "Multi-Agent System Architecture"
    title3_para = title3_frame.paragraphs[0]
    title3_para.font.size = Pt(36)
    title3_para.font.color.rgb = GOLD
    title3_para.font.bold = True
    
    # Architecture diagram
    arch_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(5.5))
    arch_frame = arch_box.text_frame
    
    architecture_text = """┌─────────────────────────────────────────────────────────────┐
│                    Web Interface Layer ✅                   │
│  ┌─────────────────┐  ┌─────────────────┐  ┌─────────────┐ │
│  │   Input Form    │  │  Results View   │  │ Goal Selector│ │
│  │  (Implemented)  │  │  (Implemented)  │  │ (Implemented)│ │
│  └─────────────────┘  └─────────────────┘  └─────────────┘ │
└─────────────────────────────────────────────────────────────┘
                                │
┌─────────────────────────────────────────────────────────────┐
│                 Agent Orchestration Layer ✅                │
│  ┌─────────────────┐  ┌─────────────────┐  ┌─────────────┐ │
│  │ Analysis Engine │  │  Goal Context   │  │ Result Cache│ │
│  │  (Implemented)  │  │  (Implemented)  │  │ (Implemented)│ │
│  └─────────────────┘  └─────────────────┘  └─────────────┘ │
└─────────────────────────────────────────────────────────────┘
                                │
┌─────────────────────────────────────────────────────────────┐
│                    Agent Processing Layer ✅                │
│  ┌─────────────┐ ┌─────────────┐ ┌─────────────┐ ┌─────────┐│
│  │Parser Agent │ │Risk Analyzer│ │Optimization │ │ Critic  ││
│  │(Implemented)│ │(Implemented)│ │(Implemented)│ │(Impl.)  ││
│  └─────────────┘ └─────────────┘ └─────────────┘ └─────────┘│
└─────────────────────────────────────────────────────────────┘"""
    
    arch_frame.text = architecture_text
    arch_para = arch_frame.paragraphs[0]
    arch_para.font.name = 'Courier New'
    arch_para.font.size = Pt(11)
    arch_para.font.color.rgb = WHITE
    
    # Continue with more slides...
    # For brevity, I'll add a few key slides. You can extend this pattern.
    
    # Slide 4: Kiro Features
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = BLUE
    bg4.line.fill.background()
    
    # Title
    title4_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title4_frame = title4_box.text_frame
    title4_frame.text = "Kiro Features Leveraged (30% of Grade)"
    title4_para = title4_frame.paragraphs[0]
    title4_para.font.size = Pt(32)
    title4_para.font.color.rgb = GOLD
    title4_para.font.bold = True
    
    # Feature grid (2x2)
    features = [
        ("Specs & Requirements", [
            "Formal requirements.md, design.md, tasks.md",
            "Property-based testing specifications", 
            "Incremental development tracking"
        ]),
        ("Agent Orchestration", [
            "4 specialized AI agents",
            "Sequential execution with data flow",
            "Error handling & graceful degradation"
        ]),
        ("Workflow Automation", [
            "Intelligent caching (2.8x performance)",
            "Goal-specific re-analysis",
            "Human-in-the-loop feedback"
        ]),
        ("Development Tools", [
            "TypeScript integration",
            "Jest + property-based testing",
            "Next.js + Mantine UI"
        ])
    ]
    
    positions = [(0.5, 1.5), (6.5, 1.5), (0.5, 4), (6.5, 4)]
    
    for i, (feature_title, feature_items) in enumerate(features):
        x, y = positions[i]
        
        # Feature box
        feature_box = slide4.shapes.add_textbox(Inches(x), Inches(y), Inches(6), Inches(2.5))
        feature_frame = feature_box.text_frame
        feature_frame.text = f"✅ {feature_title}"
        feature_para = feature_frame.paragraphs[0]
        feature_para.font.size = Pt(18)
        feature_para.font.color.rgb = GREEN
        feature_para.font.bold = True
        
        for item in feature_items:
            p = feature_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_before = Pt(4)
    
    # Add final slide: Thank You & Q&A
    final_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg_final = final_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_final.fill.solid()
    bg_final.fill.fore_color.rgb = BLUE
    bg_final.line.fill.background()
    
    # Title
    final_title_box = final_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
    final_title_frame = final_title_box.text_frame
    final_title_frame.text = "Thank You & Q&A"
    final_title_para = final_title_frame.paragraphs[0]
    final_title_para.alignment = PP_ALIGN.CENTER
    final_title_para.font.size = Pt(48)
    final_title_para.font.color.rgb = WHITE
    final_title_para.font.bold = True
    
    # Key takeaways
    takeaways_box = final_slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.33), Inches(2))
    takeaways_frame = takeaways_box.text_frame
    takeaways_frame.text = "Key Takeaways"
    takeaways_para = takeaways_frame.paragraphs[0]
    takeaways_para.alignment = PP_ALIGN.CENTER
    takeaways_para.font.size = Pt(24)
    takeaways_para.font.color.rgb = GREEN
    takeaways_para.font.bold = True
    
    takeaways = [
        "Agentic Design: Multi-agent systems excel at complex analysis",
        "Kiro Integration: Specs + workflows + testing = robust development",
        "Performance Matters: Intelligent caching transforms UX",
        "Iterative Value: Goal-based re-analysis provides perspectives"
    ]
    
    for takeaway in takeaways:
        p = takeaways_frame.add_paragraph()
        p.text = f"• {takeaway}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.CENTER
    
    # Questions prompt
    questions_box = final_slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    questions_frame = questions_box.text_frame
    questions_frame.text = "Questions & Discussion"
    questions_para = questions_frame.paragraphs[0]
    questions_para.alignment = PP_ALIGN.CENTER
    questions_para.font.size = Pt(28)
    questions_para.font.color.rgb = GOLD
    questions_para.font.bold = True
    
    # Save presentation
    prs.save('Agentic_Workflow_Reviewer_Presentation.pptx')
    print("✅ PowerPoint presentation created: Agentic_Workflow_Reviewer_Presentation.pptx")
    print("📊 Slides created: Title, Problem Statement, Architecture, Kiro Features, Thank You")
    print("🎯 Ready for Kiro Hackathon presentation!")

def main():
    """Main function"""
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not found")
        print("📦 Install with: pip install python-pptx")
        print("🔄 Then run: python create_powerpoint.py")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")

if __name__ == "__main__":
    main()